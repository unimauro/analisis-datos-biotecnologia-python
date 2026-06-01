"""
Construye el deck 'U2: ADN, ARN y Proteínas en Python' — Escuela Global.
Estética de marca: teal #158158 / cyan #24CBE5, 16:9.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

A = os.path.join(os.path.dirname(__file__), "assets")

# ---- Paleta ----
TEAL_DARK = RGBColor(0x08, 0x26, 0x2E)
TEAL      = RGBColor(0x15, 0x81, 0x58)
TEAL_MID  = RGBColor(0x0C, 0x4A, 0x50)
CYAN      = RGBColor(0x24, 0xCB, 0xE5)
INK       = RGBColor(0x0F, 0x2F, 0x33)
GREY      = RGBColor(0x5B, 0x6B, 0x6E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG   = RGBColor(0x06, 0x1E, 0x24)
CODE_FG   = RGBColor(0xE6, 0xF4, 0xF1)
CODE_KW   = RGBColor(0x36, 0xD3, 0xA6)   # keywords / acento
LIGHT     = RGBColor(0xEA, 0xF3, 0xF1)

FONT = "Calibri"
MONO = "Consolas"

EMU = 914400
SW, SH = 10.0, 5.625

prs = Presentation()
prs.slide_width = Emu(int(SW * EMU))
prs.slide_height = Emu(int(SH * EMU))
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def bg_image(s, name):
    s.shapes.add_picture(os.path.join(A, name), 0, 0, prs.slide_width, prs.slide_height)


def bg_fill(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, fnt) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = fnt
    return tb


def header(s, kicker, title):
    """Header de marca para slides de contenido."""
    rect(s, 0, 0, SW, 0.16, TEAL)                 # franja superior
    rect(s, 0, 0, 0.16, SH, TEAL_DARK)            # franja lateral marca
    txt(s, 0.55, 0.34, 9.0, 0.4,
        [[(kicker.upper(), 12, TEAL, True, FONT)]])
    txt(s, 0.55, 0.62, 9.0, 0.8,
        [[(title, 27, TEAL_DARK, True, FONT)]])
    rect(s, 0.57, 1.36, 1.4, 0.045, CYAN)         # subrayado acento


def footer(s, n):
    txt(s, 0.55, SH - 0.42, 6.0, 0.3,
        [[("Escuela Global  ·  Análisis de Datos en Biotecnología con Python  ·  U2", 9, GREY, False, FONT)]])
    txt(s, SW - 1.2, SH - 0.42, 0.6, 0.3,
        [[(str(n), 10, TEAL, True, FONT)]], align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, h, items, size=15, gap=8, color=INK):
    runs = []
    for it in items:
        if isinstance(it, tuple) and len(it) == 2 and it[1]:
            lead, rest = it
            runs.append([("●  ", size, CYAN, True, FONT), (lead, size, TEAL_DARK, True, FONT),
                         (rest, size, color, False, FONT)])
        else:
            text = it[0] if isinstance(it, tuple) else it
            runs.append([("●  ", size, CYAN, True, FONT), (text, size, color, False, FONT)])
    txt(s, x, y, w, h, runs, space_after=gap, line_spacing=1.05)


def code_box(s, x, y, w, h, lines, size=12.5, title=None):
    rect(s, x, y, w, h, CODE_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # tres "puntos" estilo editor
    for i, c in enumerate([RGBColor(0xFF,0x5F,0x56), RGBColor(0xFF,0xBD,0x2E), RGBColor(0x27,0xC9,0x3F)]):
        rect(s, x + 0.18 + i*0.20, y + 0.14, 0.11, 0.11, c, shape=MSO_SHAPE.OVAL)
    if title:
        txt(s, x + 0.95, y + 0.06, w - 1.0, 0.3, [[(title, 10, GREY, False, MONO)]])
    tb = s.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.42), Inches(w - 0.4), Inches(h - 0.55))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, kind) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.12; p.space_after = Pt(1)
        r = p.add_run(); r.text = text
        r.font.name = MONO; r.font.size = Pt(size)
        if kind == "c":      r.font.color.rgb = RGBColor(0x6E,0x8B,0x8E)  # comentario
        elif kind == "o":    r.font.color.rgb = CODE_KW                   # output/acento
        else:                r.font.color.rgb = CODE_FG
    return tb


def pic(s, name, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return s.shapes.add_picture(os.path.join(A, name), Inches(x), Inches(y), **kw)


N = 0
def num():
    global N; N += 1; return N

# ============================================================== 1. PORTADA
s = slide(); bg_image(s, "cover.png")
txt(s, 0.7, 1.15, 6.6, 0.5, [[("ESCUELA GLOBAL", 15, CYAN, True, FONT)]])
txt(s, 0.7, 1.55, 7.4, 2.0, [
    [("ADN, ARN y Proteínas", 40, WHITE, True, FONT)],
    [("en Python", 40, WHITE, True, FONT)],
], line_spacing=1.0, space_after=2)
rect(s, 0.74, 3.05, 2.0, 0.05, CYAN)
txt(s, 0.7, 3.2, 7.6, 0.9, [
    [("Unidad 2 · Representación de información biológica", 15, LIGHT, False, FONT)],
    [("Módulo: Análisis de Datos en Biotecnología Aplicada con Python", 13, RGBColor(0xB8,0xD6,0xD2), False, FONT)],
])
txt(s, 0.7, SH - 0.75, 8.0, 0.4, [[("Docente: Carlos Cárdenas Fernández", 13, WHITE, True, FONT)]])

# ============================================================== 2. AGENDA
s = slide(); bg_fill(s, WHITE); header(s, "Hoja de ruta de hoy", "Lo que veremos en esta clase")
footer(s, num())
bullets(s, 0.7, 1.7, 5.4, 3.4, [
    ("El dogma central. ", "ADN → ARN → Proteína"),
    ("ADN como datos. ", "Representarlo y medirlo en Python"),
    ("Transcripción y traducción ", "paso a paso, en código"),
    ("El código genético. ", "Codones y aminoácidos"),
    ("Biopython. ", "La librería estándar (Seq, FASTA)"),
    ("Caso en salud. ", "Detectar una mutación"),
    ("Práctica guiada ", "en Google Colab"),
], size=15, gap=10)
pic(s, "fig_helix.png", 6.7, 1.55, h=3.55)

# ============================================================== 3. DIVISOR: por qué Python
s = slide(); bg_image(s, "section.png")
txt(s, 0.9, 2.0, 8.2, 1.6, [
    [("01", 30, CYAN, True, FONT)],
    [("¿Por qué Python en biotecnología?", 32, WHITE, True, FONT)],
], space_after=6)
rect(s, 0.94, 3.55, 2.0, 0.05, CYAN)

# ============================================================== 4. POR QUÉ PYTHON
s = slide(); bg_fill(s, WHITE); header(s, "Motivación", "La biología hoy es, también, datos")
footer(s, num())
bullets(s, 0.7, 1.7, 5.6, 3.3, [
    ("Volumen. ", "Un solo genoma humano ≈ 3 200 millones de bases"),
    ("Texto = datos. ", "ADN/ARN se modelan como secuencias de letras"),
    ("Reproducibilidad. ", "Un script repite el análisis sin error humano"),
    ("Ecosistema. ", "Biopython, pandas, NumPy, Matplotlib"),
    ("Gratis y abierto. ", "Colab corre en el navegador, sin instalar"),
], size=15, gap=11)
rect(s, 6.55, 1.75, 3.0, 3.05, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 6.8, 2.0, 2.6, 2.7, [
    [("Dato", 13, TEAL, True, FONT)],
    [("3 200", 40, TEAL_DARK, True, FONT)],
    [("millones de pares de bases", 13, INK, False, FONT)],
    [("en el genoma humano", 13, INK, False, FONT)],
    [("≈ 700 MB de texto", 12, GREY, False, FONT)],
], space_after=6)

# ============================================================== 5. DOGMA CENTRAL
s = slide(); bg_fill(s, WHITE); header(s, "Concepto clave", "El Dogma Central")
footer(s, num())
pic(s, "fig_dogma.png", 0.7, 1.75, w=8.6)
bullets(s, 1.0, 4.25, 8.2, 1.0, [
    ("Transcripción: ", "el ADN se copia a ARN mensajero"),
    ("Traducción: ", "el ARN se lee en codones y se arma la proteína"),
], size=14, gap=6)

# ============================================================== 6. ADN COMO DATOS
s = slide(); bg_fill(s, WHITE); header(s, "Manos a la obra", "El ADN como un string")
footer(s, num())
code_box(s, 0.6, 1.7, 4.55, 2.7, [
    ("# Una secuencia es solo texto", "c"),
    ("adn = \"ATGCTAGCTAGCTAGC\"", "n"),
    ("", "n"),
    ("len(adn)        # longitud", "n"),
    ("adn[0]          # primera base", "n"),
    ("adn.count(\"A\")  # cuántas A", "n"),
    ("", "n"),
    (">>> 16", "o"),
    (">>> 'A'", "o"),
    (">>> 4", "o"),
], title="adn.py")
pic(s, "fig_bases.png", 5.35, 2.05, w=4.2)
txt(s, 5.4, 4.25, 4.2, 0.8,
    [[("Cada letra es una base nitrogenada. La cadena complementaria empareja A–T y G–C.",
       12.5, GREY, False, FONT)]])

# ============================================================== 7. CONTEO Y GC
s = slide(); bg_fill(s, WHITE); header(s, "Medir la secuencia", "Composición y contenido GC")
footer(s, num())
code_box(s, 0.6, 1.7, 4.7, 2.95, [
    ("def contenido_gc(seq):", "n"),
    ("    g = seq.count(\"G\")", "n"),
    ("    c = seq.count(\"C\")", "n"),
    ("    return (g + c) / len(seq) * 100", "n"),
    ("", "n"),
    ("contenido_gc(\"ATGCGCGC\")", "n"),
    (">>> 75.0", "o"),
], title="gc.py")
pic(s, "fig_gc.png", 5.5, 1.7, w=4.0)
txt(s, 0.6, 4.8, 8.8, 0.6,
    [[("El %GC indica estabilidad térmica del ADN y ayuda a comparar genes y organismos.",
       13, INK, False, FONT)]])

# ============================================================== 8. COMPLEMENTO / REVERSA
s = slide(); bg_fill(s, WHITE); header(s, "Operación esencial", "Complemento y cadena reversa")
footer(s, num())
code_box(s, 0.6, 1.7, 5.3, 3.2, [
    ("comp = {\"A\":\"T\", \"T\":\"A\",", "n"),
    ("        \"G\":\"C\", \"C\":\"G\"}", "n"),
    ("", "n"),
    ("def rev_comp(seq):", "n"),
    ("    return \"\".join(comp[b]", "n"),
    ("           for b in reversed(seq))", "n"),
    ("", "n"),
    ("rev_comp(\"ATGC\")", "n"),
    (">>> 'GCAT'", "o"),
], title="complemento.py")
bullets(s, 6.1, 1.85, 3.4, 3.0, [
    ("A ↔ T", ""),
    ("G ↔ C", ""),
    ("reversa-complementaria: ", "leemos la otra hebra 5'→3'"),
    ("Se usa para encontrar genes en ambas hebras del ADN",),
], size=14, gap=12)

# ============================================================== 9. TRANSCRIPCIÓN
s = slide(); bg_fill(s, WHITE); header(s, "Paso 1", "Transcripción: ADN → ARN")
footer(s, num())
code_box(s, 0.6, 1.75, 5.3, 2.5, [
    ("# En el ARN, la T se vuelve U", "c"),
    ("def transcribir(adn):", "n"),
    ("    return adn.replace(\"T\", \"U\")", "n"),
    ("", "n"),
    ("transcribir(\"ATGCTA\")", "n"),
    (">>> 'AUGCUA'", "o"),
], title="transcripcion.py")
rect(s, 6.2, 1.9, 3.3, 2.4, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 6.45, 2.15, 2.9, 2.0, [
    [("Regla", 13, TEAL, True, FONT)],
    [("T  →  U", 30, TEAL_DARK, True, FONT)],
    [("Timina pasa a Uracilo.", 13, INK, False, FONT)],
    [("El resto de bases no cambia.", 13, INK, False, FONT)],
], space_after=8)
txt(s, 0.6, 4.5, 8.8, 0.6,
    [[("El ARN mensajero es la copia portátil del gen que sale del núcleo hacia el ribosoma.",
       13, GREY, False, FONT)]])

# ============================================================== 10. CÓDIGO GENÉTICO
s = slide(); bg_fill(s, WHITE); header(s, "La clave de lectura", "El código genético: codones")
footer(s, num())
pic(s, "fig_codones.png", 0.7, 1.75, w=8.6)
bullets(s, 1.0, 4.2, 8.2, 1.0, [
    ("64 codones ", "(4³) codifican 20 aminoácidos + señales de inicio/STOP"),
    ("AUG ", "inicia la proteína;  UAA / UAG / UGA la terminan"),
], size=14, gap=6)

# ============================================================== 11. TRADUCCIÓN
s = slide(); bg_fill(s, WHITE); header(s, "Paso 2", "Traducción: ARN → Proteína")
footer(s, num())
code_box(s, 0.6, 1.7, 5.5, 3.3, [
    ("codigo = {\"AUG\":\"M\", \"UUU\":\"F\",", "n"),
    ("          \"GGC\":\"G\", \"UAA\":\"*\"}", "n"),
    ("", "n"),
    ("def traducir(arn):", "n"),
    ("    prot = \"\"", "n"),
    ("    for i in range(0, len(arn)-2, 3):", "n"),
    ("        prot += codigo.get(arn[i:i+3], \"?\")", "n"),
    ("    return prot", "n"),
    ("", "n"),
    ("traducir(\"AUGUUUGGC\")  >>> 'MFG'", "o"),
], title="traduccion.py", size=12)
bullets(s, 6.3, 1.85, 3.2, 3.0, [
    ("Leemos de 3 en 3 ", "(marco de lectura)"),
    ("Cada codón → 1 letra de aminoácido",),
    ("* ", "marca el fin de la proteína"),
    ("Biopython ya trae la tabla completa",),
], size=13.5, gap=11)

# ============================================================== 12. BIOPYTHON
s = slide(); bg_fill(s, WHITE); header(s, "Herramienta profesional", "Biopython: no reinventes la rueda")
footer(s, num())
code_box(s, 0.6, 1.7, 5.7, 3.35, [
    ("from Bio.Seq import Seq", "n"),
    ("", "n"),
    ("gen = Seq(\"ATGCTAGCTAGCTAA\")", "n"),
    ("gen.complement()       # complementaria", "n"),
    ("gen.reverse_complement()", "n"),
    ("gen.transcribe()       # ADN -> ARN", "n"),
    ("gen.translate()        # -> proteína", "n"),
    ("", "n"),
    (">>> Seq('MLAS*')", "o"),
], title="biopython.py", size=12)
bullets(s, 6.4, 1.85, 3.1, 3.0, [
    ("Seq ", "= secuencia con superpoderes"),
    ("SeqIO ", "lee/escribe FASTA, GenBank"),
    ("Tablas de codones ", "estándar incluidas"),
    ("pip install biopython",),
], size=13, gap=11)

# ============================================================== 13. CASO SALUD
s = slide(); bg_fill(s, WHITE); header(s, "Aplicación en salud", "Detectar una mutación")
footer(s, num())
pic(s, "fig_pipeline.png", 0.6, 1.7, w=8.8)
code_box(s, 0.9, 3.3, 8.2, 1.55, [
    ("ref  = Seq(\"ATGCGT\");  muestra = Seq(\"ATGCAT\")", "n"),
    ("dif = [i for i in range(len(ref)) if ref[i] != muestra[i]]", "n"),
    ("print(\"Mutación en posición:\", dif)   >>> Mutación en posición: [4]", "o"),
], title="mutacion.py", size=11.5)

# ============================================================== 14. BUENAS PRÁCTICAS
s = slide(); bg_fill(s, WHITE); header(s, "Evita tropiezos", "Buenas prácticas y errores comunes")
footer(s, num())
rect(s, 0.6, 1.7, 4.3, 3.2, RGBColor(0xEC,0xF7,0xF2), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 0.85, 1.9, 3.9, 0.4, [[("✓  Hazlo así", 16, TEAL, True, FONT)]])
bullets(s, 0.85, 2.4, 3.9, 2.4, [
    "Mayúsculas: ATGC, no atgc",
    "Valida que solo haya A/C/G/T",
    "Usa Biopython para producción",
    "Documenta y comenta el código",
], size=12.5, gap=8, color=INK)
rect(s, 5.1, 1.7, 4.3, 3.2, RGBColor(0xFB,0xEC,0xEC), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 5.35, 1.9, 3.9, 0.4, [[("✗  Cuidado con", 16, RGBColor(0xB2,0x3B,0x3B), True, FONT)]])
bullets(s, 5.35, 2.4, 3.9, 2.4, [
    "Olvidar el marco de lectura (de 3 en 3)",
    "Confundir T (ADN) con U (ARN)",
    "Secuencias de largo no múltiplo de 3",
    "Asumir una sola hebra del ADN",
], size=12.5, gap=8, color=INK)

# ============================================================== 15. DIVISOR PRÁCTICA
s = slide(); bg_image(s, "section.png")
txt(s, 0.9, 1.9, 8.2, 1.8, [
    [("02", 30, CYAN, True, FONT)],
    [("Práctica guiada en Google Colab", 30, WHITE, True, FONT)],
    [("Abre el notebook del repositorio y ejecútalo paso a paso", 15, LIGHT, False, FONT)],
], space_after=8)
rect(s, 0.94, 3.7, 2.0, 0.05, CYAN)

# ============================================================== 16. RETO
s = slide(); bg_fill(s, WHITE); header(s, "Tu turno", "Reto de la clase")
footer(s, num())
bullets(s, 0.7, 1.75, 8.7, 3.2, [
    ("Reto 1. ", "Calcula el %GC de una secuencia que tú elijas"),
    ("Reto 2. ", "Escribe rev_comp() y pruébalo con \"AATTCCGG\""),
    ("Reto 3. ", "Transcribe y traduce un gen hasta el primer STOP"),
    ("Reto 4. ", "Lee un archivo FASTA con Biopython e imprime su longitud"),
    ("Reto 5 (bonus). ", "Compara dos secuencias y reporta las mutaciones"),
], size=15, gap=13)
rect(s, 0.6, SH-1.05, 8.8, 0.5, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 0.85, SH-0.98, 8.4, 0.4,
    [[("Entrega: sube tu notebook .ipynb resuelto al classroom.", 12.5, TEAL_DARK, True, FONT)]])

# ============================================================== 17. RECURSOS
s = slide(); bg_fill(s, WHITE); header(s, "Para seguir aprendiendo", "Recursos y repositorio")
footer(s, num())
bullets(s, 0.7, 1.75, 8.7, 2.2, [
    ("Repositorio de la clase ", "(slides + notebook + datos)"),
    ("Biopython ", "— biopython.org / Tutorial & Cookbook"),
    ("NCBI ", "— ncbi.nlm.nih.gov (genes y secuencias reales)"),
    ("Rosalind ", "— rosalind.info (retos de bioinformática)"),
], size=14.5, gap=11)
rect(s, 0.6, 4.05, 8.8, 0.85, TEAL_DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, 0.85, 4.18, 8.4, 0.6, [
    [("github.com/unimauro/analisis-datos-biotecnologia-python", 16, CYAN, True, MONO)],
    [("Clona o abre el notebook directo en Google Colab", 12, LIGHT, False, FONT)],
], space_after=2)

# ============================================================== 18. CIERRE
s = slide(); bg_image(s, "cover.png")
txt(s, 0.7, 1.9, 8.4, 1.6, [
    [("¡Gracias!", 44, WHITE, True, FONT)],
    [("¿Preguntas?", 22, CYAN, True, FONT)],
], space_after=8)
txt(s, 0.7, SH-1.0, 8.4, 0.6, [
    [("Carlos Cárdenas Fernández  ·  Escuela Global", 14, LIGHT, True, FONT)],
    [("Análisis de Datos en Biotecnología Aplicada con Python", 12, RGBColor(0xB8,0xD6,0xD2), False, FONT)],
])

out = os.path.join(os.path.dirname(__file__), "U2_ADN_ARN_Proteinas_EscuelaGlobal.pptx")
prs.save(out)
print("OK ->", out, "|", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
